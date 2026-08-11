from dotenv import load_dotenv
from pathlib import Path
ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env')

import os
import io
import uuid
import json
import logging
from datetime import datetime, timezone, timedelta
from typing import List, Optional

import bcrypt
import jwt
import requests
from bson import ObjectId
from fastapi import FastAPI, APIRouter, Request, Response, HTTPException, Depends, UploadFile, File, Form, Header, Query
from starlette.middleware.cors import CORSMiddleware
from starlette.responses import Response as StarletteResponse
from motor.motor_asyncio import AsyncIOMotorClient
from pydantic import BaseModel, EmailStr, Field

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

mongo_url = os.environ.get('MONGO_URL', 'mongodb://localhost:27017')
db_name = os.environ.get('DB_NAME', 'studyace')
client = AsyncIOMotorClient(mongo_url)
db = client[db_name]

app = FastAPI()

# Custom Middleware for Preflight Robustness
ALLOWED_ORIGINS = {
    "https://study-ace-khaki.vercel.app",
    "https://study-ace-khaki.vercel.app/",
    "http://localhost:3000",
    "http://localhost:5173",
    "http://127.0.0.1:3000",
    "http://127.0.0.1:5173",
}

@app.middleware("http")
async def custom_cors_preflight_middleware(request: Request, call_next):
    # Handle preflight OPTIONS requests before Starlette's CORS middleware can reject them
    if request.method == "OPTIONS":
        origin = request.headers.get("origin")
        response = StarletteResponse(status_code=200)
        
        # Check explicit origins or Vercel preview URLs
        if origin in ALLOWED_ORIGINS or (origin and ".vercel.app" in origin):
            response.headers["Access-Control-Allow-Origin"] = origin
            response.headers["Access-Control-Allow-Credentials"] = "true"
            response.headers["Access-Control-Allow-Methods"] = "GET, POST, PUT, PATCH, DELETE, OPTIONS"
            response.headers["Access-Control-Allow-Headers"] = request.headers.get("access-control-request-headers", "*")
            return response
        elif origin:
            response.headers["Access-Control-Allow-Origin"] = origin
            response.headers["Access-Control-Allow-Credentials"] = "true"
            response.headers["Access-Control-Allow-Methods"] = "GET, POST, PUT, PATCH, DELETE, OPTIONS"
            response.headers["Access-Control-Allow-Headers"] = request.headers.get("access-control-request-headers", "*")
            return response
        return response
    
    response = await call_next(request)
    return response

# Standard CORS Middleware for normal requests
app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "https://study-ace-khaki.vercel.app",
        "https://study-ace-khaki.vercel.app/",
        "http://localhost:3000",
        "http://localhost:5173",
        "http://127.0.0.1:3000",
        "http://127.0.0.1:5173"
    ],
    allow_origin_regex=r"https://.*\.vercel\.app",
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
    expose_headers=["*"],
)

api_router = APIRouter(prefix="/api")

JWT_ALGORITHM = "HS256"
EMERGENT_KEY = os.environ.get("EMERGENT_LLM_KEY")
APP_NAME = "studyace"

# Object storage
STORAGE_BASE = (os.environ.get("INTEGRATION_PROXY_URL") or "").strip() or "https://integrations.emergentagent.com"
STORAGE_URL = STORAGE_BASE.strip("/") + "/objstore/api/v1/storage"
storage_key = None

def init_storage(force: bool = False):
    global storage_key
    if storage_key and not force:
        return storage_key
    if not EMERGENT_KEY:
        logger.warning("EMERGENT_LLM_KEY is not set. Storage initialization skipped.")
        return None
    try:
        resp = requests.post(f"{STORAGE_URL}/init", json={"emergent_key": EMERGENT_KEY}, timeout=30)
        resp.raise_for_status()
        storage_key = resp.json()["storage_key"]
        return storage_key
    except Exception as e:
        logger.error(f"Storage init failed: {e}")
        return None

def put_object(path: str, data: bytes, content_type: str) -> dict:
    key = init_storage()
    if not key:
        raise HTTPException(status_code=500, detail="Object storage key unavailable.")
    resp = requests.put(f"{STORAGE_URL}/objects/{path}", headers={"X-Storage-Key": key, "Content-Type": content_type}, data=data, timeout=120)
    if resp.status_code == 404:
        key = init_storage(force=True)
        resp = requests.put(f"{STORAGE_URL}/objects/{path}", headers={"X-Storage-Key": key, "Content-Type": content_type}, data=data, timeout=120)
    resp.raise_for_status()
    return resp.json()

def get_object(path: str):
    key = init_storage()
    if not key:
        raise HTTPException(status_code=500, detail="Object storage key unavailable.")
    resp = requests.get(f"{STORAGE_URL}/objects/{path}", headers={"X-Storage-Key": key}, timeout=60)
    if resp.status_code == 404:
        key = init_storage(force=True)
        resp = requests.get(f"{STORAGE_URL}/objects/{path}", headers={"X-Storage-Key": key}, timeout=60)
    resp.raise_for_status()
    return resp.content, resp.headers.get("Content-Type", "application/octet-stream")

# Text extraction
def extract_text(data: bytes, ext: str) -> str:
    ext = ext.lower()
    try:
        if ext == "pdf":
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(data))
            return "\n".join((p.extract_text() or "") for p in reader.pages)
        if ext == "docx":
            import docx
            d = docx.Document(io.BytesIO(data))
            return "\n".join(p.text for p in d.paragraphs)
        if ext == "pptx":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            out = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        out.append(shape.text_frame.text)
            return "\n".join(out)
        if ext in ("txt", "md"):
            return data.decode("utf-8", errors="ignore")
    except Exception as e:
        logger.error(f"extract_text failed: {e}")
    return ""

# Auth helpers
def hash_password(password: str) -> str:
    return bcrypt.hashpw(password.encode("utf-8"), bcrypt.gensalt()).decode("utf-8")

def verify_password(plain: str, hashed: str) -> bool:
    return bcrypt.checkpw(plain.encode("utf-8"), hashed.encode("utf-8"))

def get_jwt_secret() -> str:
    return os.environ.get("JWT_SECRET", "fallback-secret-change-me")

def create_access_token(user_id: str, email: str) -> str:
    payload = {"sub": user_id, "email": email, "exp": datetime.now(timezone.UTC) + timedelta(days=7), "type": "access"}
    return jwt.encode(payload, get_jwt_secret(), algorithm=JWT_ALGORITHM)

async def get_current_user(request: Request) -> dict:
    token = request.cookies.get("access_token")
    if not token:
        auth_header = request.headers.get("Authorization", "")
        if auth_header.startswith("Bearer "):
            token = auth_header[7:]
    if not token:
        raise HTTPException(status_code=401, detail="Not authenticated")
    try:
        payload = jwt.decode(token, get_jwt_secret(), algorithms=[JWT_ALGORITHM])
        user = await db.users.find_one({"_id": ObjectId(payload["sub"])})
        if not user:
            raise HTTPException(status_code=401, detail="User not found")
        user["id"] = str(user["_id"])
        user.pop("_id", None)
        user.pop("password_hash", None)
        return user
    except jwt.ExpiredSignatureError:
        raise HTTPException(status_code=401, detail="Token expired")
    except jwt.InvalidTokenError:
        raise HTTPException(status_code=401, detail="Invalid token")

def set_auth_cookie(response: Response, token: str):
    response.set_cookie(key="access_token", value=token, httponly=True, secure=True, samesite="none", max_age=604800, path="/")

# Models
class RegisterInput(BaseModel):
    name: str
    email: EmailStr
    password: str

class LoginInput(BaseModel):
    email: EmailStr
    password: str

class CourseInput(BaseModel):
    name: str
    color: Optional[str] = "#8B5CF6"
    description: Optional[str] = ""

class CourseUpdateInput(BaseModel):
    name: Optional[str] = None
    color: Optional[str] = None
    description: Optional[str] = None
    archived: Optional[bool] = None

class QuizGenInput(BaseModel):
    document_ids: List[str]
    quiz_type: str
    num_questions: int = 10
    title: Optional[str] = None
    topics: Optional[str] = ""
    folder_id: Optional[str] = None

class FolderInput(BaseModel):
    name: str

class QuizUpdateInput(BaseModel):
    title: Optional[str] = None
    folder_id: Optional[str] = None

class QuizAttemptInput(BaseModel):
    answers: List[str]

class ChatInput(BaseModel):
    document_ids: List[str] = []
    message: str
    session_id: Optional[str] = None

class StudyPlanInput(BaseModel):
    exam_name: str
    exam_date: str
    daily_hours: float
    topics: str = ""
    timetable: str = ""
    syllabus: str = ""

def now_iso():
    return datetime.now(timezone.UTC).isoformat()

# LLM
FORMAT_RULES = """FORMATTING RULES (very important, follow exactly):
- Write clean, readable prose in Markdown. Use ## for main headings, ### for sub-headings, and **bold** for key terms/labels. Use short paragraphs, bullet lists, and leave a blank line between sections. Maintain clear visual hierarchy.
- NEVER use LaTeX, dollar signs ($), or backslash math commands (no \\frac, \\times, \\div, etc.). NEVER output raw #, ##, ### as literal text mid-sentence.
- Use REAL Unicode math symbols: × (multiply), ÷ (divide), − (minus), ± ≤ ≥ ≠ ≈ √ ∑ ∫ ∆ π θ ° → ∞, and superscripts/subscripts like x², xⁿ, H₂O.
- For a stacked fraction (numerator on top, denominator on bottom) write EXACTLY [[frac:NUMERATOR|DENOMINATOR]] — e.g. [[frac:a + b|2c]]. Write ONE well-formed token per fraction: never add extra brackets and never put text after the closing ]]. WRONG: [[frac:1|2]]mv²]] → RIGHT: [[frac:mv²|2]] (or ½mv²). WRONG: [[frac:∆s|∆t]]]] → RIGHT: [[frac:∆s|∆t]]. For simple values you may use ½, ¾. Put display equations centered on their own line using [[center:...]].
- Keep the screen uncluttered and the text naturally formatted, like a great textbook."""
ATTACH_EXTS = ("pdf", "png", "jpg", "jpeg", "webp", "gif")

async def fetch_doc_materials(doc_ids, user_id, attach_files=True, per_text=45000, max_files=6, max_bytes=18 * 1024 * 1024):
    docs = await db.documents.find({"id": {"$in": doc_ids}, "user_id": user_id, "is_deleted": False}, {"_id": 0}).to_list(100)
    material = ""
    files = []
    total = 0
    for d in docs:
        t = (d.get("text") or "")[:per_text]
        material += f"\n\n=== DOCUMENT: {d['original_filename']} ===\n{t}"
        if attach_files and d.get("ext") in ATTACH_EXTS and len(files) < max_files:
            try:
                from emergentintegrations.llm.chat import FileContent
                content, ctype = get_object(d["storage_path"])
                if total + len(content) <= max_bytes:
                    import base64 as _b64
                    mime = d.get("content_type") or ctype or "application/pdf"
                    files.append(FileContent(mime, _b64.b64encode(content).decode("utf-8")))
                    total += len(content)
            except Exception as e:
                logger.error(f"attach file failed: {e}")
    return material[:160000], files, docs

async def llm_generate(system_message: str, prompt: str, session_id: str = None, file_contents=None) -> str:
    from emergentintegrations.llm.chat import LlmChat, UserMessage
    chat = LlmChat(api_key=EMERGENT_KEY, session_id=session_id or str(uuid.uuid4()), system_message=system_message).with_model("gemini", "gemini-3-flash-preview")
    resp = await chat.send_message(UserMessage(text=prompt, file_contents=file_contents or []))
    return resp

def parse_json_block(text: str):
    t = text.strip()
    if t.startswith("```"):
        t = t.split("```", 2)[1]
        if t.startswith("json"):
            t = t[4:]
    t = t.strip().strip("`").strip()
    start = t.find("{")
    if start == -1:
        start = t.find("[")
    end = max(t.rfind("}"), t.rfind("]"))
    if start != -1 and end != -1:
        t = t[start:end + 1]
    return json.loads(t)

import re as _re
def norm_frac(s):
    if not isinstance(s, str):
        return s
    return _re.sub(r'\]{3,}', ']]', s)

def sanitize_obj(o):
    if isinstance(o, str):
        return norm_frac(o)
    if isinstance(o, list):
        return [sanitize_obj(x) for x in o]
    if isinstance(o, dict):
        return {k: sanitize_obj(v) for k, v in o.items()}
    return o

# -------------------- Auth routes --------------------
@api_router.post("/auth/register")
async def register(data: RegisterInput, response: Response):
    email = data.email.lower()
    if await db.users.find_one({"email": email}):
        raise HTTPException(status_code=400, detail="Email already registered")
    doc = {"name": data.name, "email": email, "password_hash": hash_password(data.password), "role": "user", "created_at": now_iso(), "prefs": {"daily_hours": 2}}
    res = await db.users.insert_one(doc)
    uid = str(res.insert_id)
    set_auth_cookie(response, create_access_token(uid, email))
    return {"id": uid, "name": data.name, "email": email, "role": "user"}

@api_router.post("/auth/login")
async def login(data: LoginInput, response: Response):
    email = data.email.lower()
    user = await db.users.find_one({"email": email})
    if not user or not verify_password(data.password, user["password_hash"]):
        raise HTTPException(status_code=401, detail="Invalid email or password")
    uid = str(user["_id"])
    set_auth_cookie(response, create_access_token(uid, email))
    return {"id": uid, "name": user["name"], "email": email, "role": user.get("role", "user")}

@api_router.post("/auth/logout")
async def logout(response: Response, user: dict = Depends(get_current_user)):
    response.delete_cookie("access_token", path="/")
    return {"ok": True}

@api_router.get("/auth/me")
async def me(user: dict = Depends(get_current_user)):
    return user

# -------------------- Course routes --------------------
@api_router.post("/courses")
async def create_course(data: CourseInput, user: dict = Depends(get_current_user)):
    doc = {"id": str(uuid.uuid4()), "user_id": user["id"], "name": data.name, "color": data.color, "description": data.description, "archived": False, "created_at": now_iso()}
    await db.courses.insert_one(doc)
    doc.pop("_id", None)
    return doc

@api_router.get("/courses")
async def list_courses(user: dict = Depends(get_current_user)):
    courses = await db.courses.find({"user_id": user["id"]}, {"_id": 0}).sort("created_at", -1).to_list(1000)
    for c in courses:
        c["archived"] = c.get("archived", False)
        c["doc_count"] = await db.documents.count_documents({"course_id": c["id"], "is_deleted": False})
        c["quiz_count"] = await db.quizzes.count_documents({"course_id": c["id"]})
    return courses

@api_router.get("/courses/{course_id}")
async def get_course(course_id: str, user: dict = Depends(get_current_user)):
    c = await db.courses.find_one({"id": course_id, "user_id": user["id"]}, {"_id": 0})
    if not c:
        raise HTTPException(status_code=404, detail="Course not found")
    return c

@api_router.patch("/courses/{course_id}")
async def update_course(course_id: str, data: CourseUpdateInput, user: dict = Depends(get_current_user)):
    update = {k: v for k, v in data.model_dump(exclude_unset=True).items()}
    if not update:
        raise HTTPException(status_code=400, detail="Nothing to update")
    res = await db.courses.update_one({"id": course_id, "user_id": user["id"]}, {"$set": update})
    if res.matched_count == 0:
        raise HTTPException(status_code=404, detail="Course not found")
    return {"ok": True}

@api_router.delete("/courses/{course_id}")
async def delete_course(course_id: str, user: dict = Depends(get_current_user)):
    await db.courses.delete_one({"id": course_id, "user_id": user["id"]})
    await db.documents.update_many({"course_id": course_id}, {"$set": {"is_deleted": True}})
    await db.quizzes.delete_many({"course_id": course_id})
    await db.folders.delete_many({"course_id": course_id})
    return {"ok": True}

# -------------------- Document routes --------------------
async def suggest_doc_name(text, filename, ext, data=None, ctype=None):
    try:
        files = []
        if (not text or len(text.strip()) < 40) and ext in ATTACH_EXTS and data:
            from emergentintegrations.llm.chat import FileContent
            import base64 as _b64
            files = [FileContent(ctype or "application/pdf", _b64.b64encode(data).decode("utf-8"))]
        prompt = f"Suggest a concise, descriptive title (3 to 8 words, Title Case, NO file extension, no quotes) for this study document based on its actual content. Return ONLY the title.\n\nOriginal filename: {filename}\nContent excerpt:\n{text or ''}[:4000]"
        name = await llm_generate("You name study documents concisely and accurately.", prompt, file_contents=files)
        name = (name or "").strip().strip('"').split("\n")[0][:80]
        return name or filename
    except Exception as e:
        logger.error(f"name suggest failed: {e}")
        return filename

@api_router.post("/courses/{course_id}/documents")
async def upload_document(course_id: str, file: UploadFile = File(...), user: dict = Depends(get_current_user)):
    course = await db.courses.find_one({"id": course_id, "user_id": user["id"]})
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")
    ext = file.filename.split(".")[-1].lower() if "." in file.filename else "bin"
    data = await file.read()
    path = f"{APP_NAME}/uploads/{user['id']}/{uuid.uuid4()}.{ext}"
    result = put_object(path, data, file.content_type or "application/octet-stream")
    text = extract_text(data, ext)
    suggested = await suggest_doc_name(text, file.filename, ext, data, file.content_type)
    doc = {"id": str(uuid.uuid4()), "user_id": user["id"], "course_id": course_id, "storage_path": result["path"], "original_filename": file.filename, "name": suggested, "suggested_name": suggested, "content_type": file.content_type, "ext": ext, "size": result.get("size", len(data)), "text": text[:200000], "is_deleted": False, "created_at": now_iso()}
    await db.documents.insert_one(doc)
    return {k: v for k, v in doc.items() if k not in {"_id", "text"}}

@api_router.get("/courses/{course_id}/documents")
async def list_documents(course_id: str, user: dict = Depends(get_current_user)):
    docs = await db.documents.find({"course_id": course_id, "user_id": user["id"], "is_deleted": False}, {"_id": 0, "text": 0}).sort("created_at", -1).to_list(1000)
    return docs

@api_router.get("/documents/{doc_id}/text")
async def get_document_text(doc_id: str, user: dict = Depends(get_current_user)):
    d = await db.documents.find_one({"id": doc_id, "user_id": user["id"]}, {"_id": 0})
    if not d:
        raise HTTPException(status_code=404, detail="Document not found")
    return {"id": d["id"], "filename": d["original_filename"], "text": d.get("text", ""), "ext": d["ext"]}

@api_router.get("/documents/{doc_id}/download")
async def download_document(doc_id: str, request: Request, authorization: str = Header(None), auth: str = Query(None)):
    token = request.cookies.get("access_token")
    if not token and authorization and authorization.startswith("Bearer "):
        token = authorization[7:]
    elif not token and auth:
        token = auth
    if not token:
        raise HTTPException(status_code=401, detail="Not authenticated")
    try:
        payload = jwt.decode(token, get_jwt_secret(), algorithms=[JWT_ALGORITHM])
        user_id = payload["sub"]
    except jwt.InvalidTokenError:
        raise HTTPException(status_code=401, detail="Invalid token")
    d = await db.documents.find_one({"id": doc_id, "user_id": user_id, "is_deleted": False})
    if not d:
        raise HTTPException(status_code=404, detail="Document not found")
    content, ctype = get_object(d["storage_path"])
    return StarletteResponse(content=content, media_type=d.get("content_type") or ctype, headers={"Content-Disposition": f'inline; filename="{d["original_filename"]}"'})

@api_router.delete("/documents/{doc_id}")
async def delete_document(doc_id: str, user: dict = Depends(get_current_user)):
    await db.documents.update_one({"id": doc_id, "user_id": user["id"]}, {"$set": {"is_deleted": True}})
    return {"ok": True}

# -------------------- Quizzes --------------------
QUIZ_INSTRUCTIONS = {
    "mcq": "Each question object: {\"question\": str, \"options\": [4 strings], \"answer\": str (must exactly match one option), \"explanation\": str, \"source\": str (short quote or section reference from the material)}",
    "flashcards": "Each question object: {\"question\": str (front/term), \"answer\": str (back/definition), \"explanation\": str, \"source\": str). Leave options as empty list.",
    "short_answer": "Each question object: {\"question\": str, \"answer\": str (concise ideal answer), \"explanation\": str, \"source\": str). Leave options as empty list.",
    "fill_blank": "Each question object: {\"question\": str (sentence with a __ blank), \"answer\": str (the word/phrase for the blank), \"explanation\": str, \"source\": str). Leave options as empty list.",
}

@api_router.post("/courses/{course_id}/quizzes/generate")
async def generate_quiz(course_id: str, data: QuizGenInput, user: dict = Depends(get_current_user)):
    course = await db.courses.find_one({"id": course_id, "user_id": user["id"]})
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")
    docs = await db.documents.find({"id": {"$in": data.document_ids}, "user_id": user["id"]}, {"_id": 0}).to_list(100)
    if not docs:
        raise HTTPException(status_code=400, detail="No documents selected")
    material, file_contents, _ = await fetch_doc_materials(data.document_ids, user["id"])
    num_q = max(1, min(50, int(data.num_questions)))
    topics_line = f"Focus specifically on these topics: {data.topics}." if data.topics else ""
    system = "You are an expert tutor that creates high-quality study assessments. You always respond with VALID JSON only (no markdown fences, no prose)."
    prompt = f"""Create a {data.quiz_type} quiz with exactly {num_q} questions from the study material below (documents may be attached as files — read their words, formulas, diagrams and pictures too). {topics_line}
{QUIZ_INSTRUCTIONS.get(data.quiz_type, QUIZ_INSTRUCTIONS['mcq'])}
IMPORTANT MIX: Do NOT make every question pure theory. When the material contains formulas, equations, or worked examples, include a healthy mix of formula-based / calculation / apply-the-formula questions alongside conceptual/theory questions.
MATH FORMATTING inside any string: use real Unicode symbols (× ÷ − ± ≤ ≥ ≠ √ π θ ° x²), NEVER LaTeX, $ or backslash commands. For stacked fractions write [[frac:numerator|denominator]].
For every question, 'source' MUST reference the specific document name and page/slide/section the answer is drawn from.
Respond ONLY with JSON: {{"questions": [ ...question objects... ]}}
STUDY MATERIAL:
{material}"""
    try:
        raw = await llm_generate(system, prompt, file_contents=file_contents)
        parsed = parse_json_block(raw)
        questions = parsed.get("questions", parsed) if isinstance(parsed, dict) else parsed
        questions = sanitize_obj(questions)
    except Exception as e:
        logger.error(f"quiz gen failed: {e}")
        raise HTTPException(status_code=500, detail="AI quiz generation failed. Please try again.")
    quiz = {"id": str(uuid.uuid4()), "user_id": user["id"], "course_id": course_id, "title": data.title or f"{data.quiz_type.replace('_',' ').title()} Quiz", "quiz_type": data.quiz_type, "num_questions": len(questions), "document_ids": data.document_ids, "topics": data.topics, "folder_id": data.folder_id, "questions": questions, "attempts": [], "best_score": None, "created_at": now_iso()}
    await db.quizzes.insert_one(quiz)
    quiz.pop("_id", None)
    return quiz

@api_router.get("/courses/{course_id}/quizzes")
async def list_quizzes(course_id: str, user: dict = Depends(get_current_user)):
    quizzes = await db.quizzes.find({"course_id": course_id, "user_id": user["id"]}, {"_id": 0, "questions": 0}).sort("created_at", -1).to_list(1000)
    return quizzes

@api_router.get("/quizzes/{quiz_id}")
async def get_quiz(quiz_id: str, user: dict = Depends(get_current_user)):
    q = await db.quizzes.find_one({"id": quiz_id, "user_id": user["id"]}, {"_id": 0})
    if not q:
        raise HTTPException(status_code=404, detail="Quiz not found")
    return q

@api_router.patch("/quizzes/{quiz_id}")
async def update_quiz(quiz_id: str, data: QuizUpdateInput, user: dict = Depends(get_current_user)):
    update = {k: v for k, v in data.model_dump(exclude_unset=True).items()}
    if not update:
        raise HTTPException(status_code=400, detail="Nothing to update")
    res = await db.quizzes.update_one({"id": quiz_id, "user_id": user["id"]}, {"$set": update})
    if res.matched_count == 0:
        raise HTTPException(status_code=404, detail="Quiz not found")
    return {"ok": True}

@api_router.delete("/quizzes/{quiz_id}")
async def delete_quiz(quiz_id: str, user: dict = Depends(get_current_user)):
    await db.quizzes.delete_one({"id": quiz_id, "user_id": user["id"]})
    return {"ok": True}

# -------------------- Folder routes --------------------
@api_router.post("/courses/{course_id}/folders")
async def create_folder(course_id: str, data: FolderInput, user: dict = Depends(get_current_user)):
    course = await db.courses.find_one({"id": course_id, "user_id": user["id"]})
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")
    doc = {"id": str(uuid.uuid4()), "user_id": user["id"], "course_id": course_id, "name": data.name, "created_at": now_iso()}
    await db.folders.insert_one(doc)
    doc.pop("_id", None)
    return doc

@api_router.get("/courses/{course_id}/folders")
async def list_folders(course_id: str, user: dict = Depends(get_current_user)):
    folders = await db.folders.find({"course_id": course_id, "user_id": user["id"]}, {"_id": 0}).sort("created_at", 1).to_list(1000)
    for f in folders:
        f["quiz_count"] = await db.quizzes.count_documents({"folder_id": f["id"]})
    return folders

@api_router.delete("/folders/{folder_id}")
async def delete_folder(folder_id: str, user: dict = Depends(get_current_user)):
    await db.folders.delete_one({"id": folder_id, "user_id": user["id"]})
    await db.quizzes.update_many({"folder_id": folder_id, "user_id": user["id"]}, {"$set": {"folder_id": None}})
    return {"ok": True}

def normalize(s: str) -> str:
    return ''.join(ch.lower() for ch in (s or "") if ch.isalnum())

@api_router.post("/quizzes/{quiz_id}/attempt")
async def submit_attempt(quiz_id: str, data: QuizAttemptInput, user: dict = Depends(get_current_user)):
    q = await db.quizzes.find_one({"id": quiz_id, "user_id": user["id"]})
    if not q:
        raise HTTPException(status_code=404, detail="Quiz not found")
    questions = q["questions"]
    results = []
    correct = 0
    for i, ques in enumerate(questions):
        user_ans = data.answers[i] if i < len(data.answers) else ""
        gold = ques.get("answer", "")
        if q["quiz_type"] == "mcq":
            is_correct = normalize(user_ans) == normalize(gold)
        else:
            na, ng = normalize(user_ans), normalize(gold)
            is_correct = bool(na) and (na == ng or (len(ng) > 3 and (na in ng or ng in na)))
        if is_correct:
            correct += 1
        results.append({"question": ques.get("question"), "user_answer": user_ans, "correct_answer": gold, "is_correct": is_correct, "explanation": ques.get("explanation", ""), "source": ques.get("source", ""), "options": ques.get("options", [])})
    score = round(100 * correct / len(questions)) if questions else 0
    attempt = {"score": score, "correct": correct, "total": len(questions), "at": now_iso(), "weak": [r["question"] for r in results if not r["is_correct"]]}
    best = q.get("best_score")
    new_best = score if best is None else max(best, score)
    await db.quizzes.update_one({"id": quiz_id}, {"$push": {"attempts": attempt}, "$set": {"best_score": new_best}})
    return {"score": score, "correct": correct, "total": len(questions), "results": results}

# -------------------- Chat routes --------------------
@api_router.post("/courses/{course_id}/chat")
async def chat_with_docs(course_id: str, data: ChatInput, user: dict = Depends(get_current_user)):
    session_id = data.session_id or str(uuid.uuid4())
    context = ""
    if data.document_ids:
        docs = await db.documents.find({"id": {"$in": data.document_ids}, "user_id": user["id"]}, {"_id": 0}).to_list(100)
        for d in docs:
            context += f"\n\n== {d['original_filename']} ==\n{d.get('text', '')[:30000]}"
    history = await db.chat_messages.find({"session_id": session_id, "user_id": user["id"]}, {"_id": 0}).sort("created_at", 1).to_list(50)
    hist_text = "\n".join(f"{m['role']}: {m['content']}" for m in history[-10:])
    system = "You are StudyAce, a friendly, encouraging AI tutor. Answer using the provided study material. Cite the document name when referencing facts. If the answer isn't in the material, say so and give your best general guidance. Keep answers focused and well-structured.\n\n" + FORMAT_RULES
    prompt = f"STUDY MATERIAL:{context[:100000] or ' (none selected)'}\n\nCONVERSATION:\n{hist_text}\n\nStudent: {data.message}\nStudyAce:"
    try:
        answer = await llm_generate(system, prompt, session_id)
    except Exception as e:
        logger.error(f"chat failed: {e}")
        raise HTTPException(status_code=500, detail="AI chat failed. Please try again.")
    await db.chat_messages.insert_one({"session_id": session_id, "user_id": user["id"], "course_id": course_id, "role": "user", "content": data.message, "created_at": now_iso()})
    answer = norm_frac(answer)
    await db.chat_messages.insert_one({"session_id": session_id, "user_id": user["id"], "course_id": course_id, "role": "assistant", "content": answer, "created_at": now_iso()})
    return {"session_id": session_id, "answer": answer}

@api_router.get("/courses/{course_id}/chat/{session_id}")
async def get_chat_history(course_id: str, session_id: str, user: dict = Depends(get_current_user)):
    msgs = await db.chat_messages.find({"session_id": session_id, "user_id": user["id"]}, {"_id": 0}).sort("created_at", 1).to_list(200)
    return msgs

# -------------------- Study plan --------------------
@api_router.post("/courses/{course_id}/study-plan")
async def create_study_plan(course_id: str, data: StudyPlanInput, user: dict = Depends(get_current_user)):
    course = await db.courses.find_one({"id": course_id, "user_id": user["id"]})
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")
    system = "You are an expert academic study planner. You respond ONLY with valid JSON, no markdown."
    prompt = f"""Create a detailed day-by-day study plan for a student.
Course: {course['name']}
Exam: {data.exam_name} on {data.exam_date}
Available study hours per day: {data.daily_hours}
Weekly timetable / busy times: {data.timetable or 'not provided'}
Syllabus / topics to cover: {data.syllabus or data.topics or 'not provided'}
Today's date: {datetime.now(timezone.utc).date().isoformat()}
Distribute topics sensibly across days leading up to the exam, leaving the last day for revision. Respect busy times.
Respond ONLY with JSON:
{{"overview": "1-2 sentence summary", "days": [{{"date": "YYYY-MM-DD", "focus": "topic focus", "tasks": ["task1","task2"], "hours": number}}], "tips": ["tip1","tip2"]}}"""
    try:
        raw = await llm_generate(system, prompt)
        plan = parse_json_block(raw)
    except Exception as e:
        logger.error(f"study plan failed: {e}")
        raise HTTPException(status_code=500, detail="AI study plan generation failed. Please try again.")
    doc = {"id": str(uuid.uuid4()), "user_id": user["id"], "course_id": course_id, "exam_name": data.exam_name, "exam_date": data.exam_date, "daily_hours": data.daily_hours, "plan": plan, "created_at": now_iso()}
    await db.study_plans.insert_one(doc)
    doc.pop("_id", None)
    return doc

@api_router.get("/courses/{course_id}/study-plans")
async def list_study_plans(course_id: str, user: dict = Depends(get_current_user)):
    plans = await db.study_plans.find({"course_id": course_id, "user_id": user["id"]}, {"_id": 0}).sort("created_at", -1).to_list(100)
    return plans

@api_router.delete("/study-plans/{plan_id}")
async def delete_plan(plan_id: str, user: dict = Depends(get_current_user)):
    await db.study_plans.delete_one({"id": plan_id, "user_id": user["id"]})
    return {"ok": True}

@api_router.post("/courses/{course_id}/schedule-file")
async def parse_schedule_file(course_id: str, kind: str = Form("syllabus"), file: UploadFile = File(...), user: dict = Depends(get_current_user)):
    course = await db.courses.find_one({"id": course_id, "user_id": user["id"]})
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")
    ext = file.filename.split(".")[-1].lower() if "." in file.filename else "bin"
    data = await file.read()
    path = f"{APP_NAME}/uploads/{user['id']}/{uuid.uuid4()}.{ext}"
    result = put_object(path, data, file.content_type or "application/octet-stream")
    text = extract_text(data, ext)
    files = []
    if ext in ATTACH_EXTS:
        try:
            from emergentintegrations.llm.chat import FileContent
            import base64 as _b64
            files = [FileContent(file.content_type or "application/pdf", _b64.b64encode(data).decode("utf-8"))]
        except Exception as e:
            logger.error(f"Schedule file attachment failed: {e}")
    system = "You are an expert academic organizer. Extract structured information from study schedules, timetables, or syllabi."
    prompt = f"Extract all topics, key dates, or timetable schedules from this uploaded {kind} file: \n\n{text[:100000]}"
    try:
        extracted = await llm_generate(system, prompt, file_contents=files)
    except Exception as e:
        logger.error(f"parse schedule file failed: {e}")
        extracted = text[:2000]
    return {"filename": file.filename, "text": extracted}

# -------------------- Dashboard --------------------
@api_router.get("/dashboard")
async def get_dashboard(user: dict = Depends(get_current_user)):
    courses = await db.courses.find({"user_id": user["id"], "archived": {"$ne": True}}).to_list(1000)
    course_ids = [c["id"] for c in courses]
    doc_count = await db.documents.count_documents({"user_id": user["id"], "course_id": {"$in": course_ids}, "is_deleted": False})
    quiz_count = await db.quizzes.count_documents({"user_id": user["id"], "course_id": {"$in": course_ids}})
    attempts = []
    async for q in db.quizzes.find({"user_id": user["id"], "course_id": {"$in": course_ids}}, {"attempts": 1}):
        attempts.extend(q.get("attempts", []))
    avg_score = round(sum(a.get("score", 0) for a in attempts) / len(attempts)) if attempts else None
    weak_areas = []
    for a in attempts:
        for w in a.get("weak", []):
            weak_areas.append({"question": w, "misses": 1})
    upcoming_exams = await db.study_plans.find({"user_id": user["id"], "course_id": {"$in": course_ids}}, {"_id": 0, "exam_name": 1, "exam_date": 1}).sort("exam_date", 1).to_list(5)
    return {
        "courses": len(courses),
        "documents": doc_count,
        "quizzes": quiz_count,
        "attempts": len(attempts),
        "avg_score": avg_score,
        "weak_areas": weak_areas[:10],
        "upcoming_exams": upcoming_exams,
        "recent_quizzes": []
    }

# -------------------- Walkthrough --------------------
@api_router.post("/courses/{course_id}/walkthrough/generate")
async def generate_walkthrough(course_id: str, data: dict, user: dict = Depends(get_current_user)):
    course = await db.courses.find_one({"id": course_id, "user_id": user["id"]})
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")
    doc_ids = data.get("document_ids", [])
    if not doc_ids:
        raise HTTPException(status_code=400, detail="No documents selected")
    material, file_contents, docs = await fetch_doc_materials(doc_ids, user["id"])
    system = "You are an expert educator creating an interactive walkthrough. Respond with VALID JSON only."
    prompt = f"""Create an interactive step-by-step walkthrough from the following study material. Each step should be one of: concept, example, realworld, or question. Include a heading, clear content, and a source reference.
Format as JSON: {{"steps": [{{"type": "concept|example|realworld|question", "heading": "Step title", "content": "explanation", "source": "doc name", "question": "optional question", "answer": "optional answer"}}]}}
Title: {data.get('title') or f'Walkthrough: {course["name"]}'}
Material:
{material[:120000]}"""
    try:
        raw = await llm_generate(system, prompt, file_contents=file_contents)
        parsed = parse_json_block(raw)
        steps = parsed.get("steps", parsed) if isinstance(parsed, dict) else parsed
        steps = sanitize_obj(steps)
    except Exception as e:
        logger.error(f"walkthrough gen failed: {e}")
        raise HTTPException(status_code=500, detail="Walkthrough generation failed. Please try again.")
    doc = {"id": str(uuid.uuid4()), "user_id": user["id"], "course_id": course_id, "title": data.get('title') or f'Walkthrough: {course["name"]}', "steps": steps, "progress": 0, "total_steps": len(steps), "created_at": now_iso()}
    await db.walkthroughs.insert_one(doc)
    doc.pop("_id", None)
    return doc

@api_router.get("/courses/{course_id}/walkthroughs")
async def list_walkthroughs(course_id: str, user: dict = Depends(get_current_user)):
    items = await db.walkthroughs.find({"course_id": course_id, "user_id": user["id"]}, {"_id": 0}).sort("created_at", -1).to_list(100)
    return items

@api_router.get("/walkthroughs/{wid}")
async def get_walkthrough(wid: str, user: dict = Depends(get_current_user)):
    w = await db.walkthroughs.find_one({"id": wid, "user_id": user["id"]}, {"_id": 0})
    if not w:
        raise HTTPException(status_code=404, detail="Walkthrough not found")
    return w

@api_router.patch("/walkthroughs/{wid}/progress")
async def update_walkthrough_progress(wid: str, data: dict, user: dict = Depends(get_current_user)):
    current_index = data.get("current_index", 0)
    res = await db.walkthroughs.update_one({"id": wid, "user_id": user["id"]}, {"$set": {"progress": current_index}})
    if res.matched_count == 0:
        raise HTTPException(status_code=404, detail="Walkthrough not found")
    return {"ok": True}

@api_router.delete("/walkthroughs/{wid}")
async def delete_walkthrough(wid: str, user: dict = Depends(get_current_user)):
    await db.walkthroughs.delete_one({"id": wid, "user_id": user["id"]})
    return {"ok": True}

# -------------------- Study Guide --------------------
@api_router.post("/courses/{course_id}/study-guide/generate")
async def generate_study_guide(course_id: str, data: dict, user: dict = Depends(get_current_user)):
    course = await db.courses.find_one({"id": course_id, "user_id": user["id"]})
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")
    doc_ids = data.get("document_ids", [])
    if not doc_ids:
        raise HTTPException(status_code=400, detail="No documents selected")
    material, file_contents, docs = await fetch_doc_materials(doc_ids, user["id"])
    system = "You are an expert study guide creator. Respond with VALID JSON only."
    prompt = f"""Create a comprehensive study guide from the following material. Include all key concepts, definitions, and important details. Format as markdown with proper headings.
Format as JSON: {{"content": "markdown content here"}}
Title: {data.get('title') or f'Study Guide: {course["name"]}'}
Material:
{material[:120000]}"""
    try:
        raw = await llm_generate(system, prompt, file_contents=file_contents)
        parsed = parse_json_block(raw)
        content = parsed.get("content", parsed) if isinstance(parsed, dict) else raw
    except Exception as e:
        logger.error(f"study guide gen failed: {e}")
        raise HTTPException(status_code=500, detail="Study guide generation failed. Please try again.")
    doc = {"id": str(uuid.uuid4()), "user_id": user["id"], "course_id": course_id, "title": data.get('title') or f'Study Guide: {course["name"]}', "content": sanitize_obj(content), "checklist_state": {}, "created_at": now_iso()}
    await db.study_guides.insert_one(doc)
    doc.pop("_id", None)
    return doc

@api_router.get("/courses/{course_id}/study-guides")
async def list_study_guides(course_id: str, user: dict = Depends(get_current_user)):
    items = await db.study_guides.find({"course_id": course_id, "user_id": user["id"]}, {"_id": 0}).sort("created_at", -1).to_list(100)
    return items

@api_router.get("/study-guides/{gid}")
async def get_study_guide(gid: str, user: dict = Depends(get_current_user)):
    g = await db.study_guides.find_one({"id": gid, "user_id": user["id"]}, {"_id": 0})
    if not g:
        raise HTTPException(status_code=404, detail="Study guide not found")
    return g

@api_router.patch("/study-guides/{gid}/checklist")
async def update_checklist(gid: str, data: dict, user: dict = Depends(get_current_user)):
    res = await db.study_guides.update_one({"id": gid, "user_id": user["id"]}, {"$set": {"checklist_state": data.get("checklist_state", {})}})
    if res.matched_count == 0:
        raise HTTPException(status_code=404, detail="Study guide not found")
    return {"ok": True}

@api_router.delete("/study-guides/{gid}")
async def delete_study_guide(gid: str, user: dict = Depends(get_current_user)):
    await db.study_guides.delete_one({"id": gid, "user_id": user["id"]})
    return {"ok": True}

# -------------------- Key Terms --------------------
@api_router.post("/courses/{course_id}/key-terms/generate")
async def generate_key_terms(course_id: str, data: dict, user: dict = Depends(get_current_user)):
    course = await db.courses.find_one({"id": course_id, "user_id": user["id"]})
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")
    doc_ids = data.get("document_ids", [])
    if not doc_ids:
        raise HTTPException(status_code=400, detail="No documents selected")
    material, file_contents, docs = await fetch_doc_materials(doc_ids, user["id"])
    system = "You are an expert at extracting key terms. Respond with VALID JSON only."
    prompt = f"""Extract the most important key terms from the following material. For each term, provide a clear definition and the source document.
Format as JSON: {{"terms": [{{"term": "term name", "definition": "definition", "source": "doc name"}}]}}
Title: {data.get('title') or f'Key Terms: {course["name"]}'}
Material:
{material[:120000]}"""
    try:
        raw = await llm_generate(system, prompt, file_contents=file_contents)
        parsed = parse_json_block(raw)
        terms = parsed.get("terms", parsed) if isinstance(parsed, dict) else parsed
        terms = sanitize_obj(terms)
        terms = sorted(terms, key=lambda x: x.get("term", "").lower())
    except Exception as e:
        logger.error(f"key terms gen failed: {e}")
        raise HTTPException(status_code=500, detail="Key terms generation failed. Please try again.")
    doc = {"id": str(uuid.uuid4()), "user_id": user["id"], "course_id": course_id, "title": data.get('title') or f'Key Terms: {course["name"]}', "terms": terms, "created_at": now_iso()}
    await db.key_terms.insert_one(doc)
    doc.pop("_id", None)
    return doc

@api_router.get("/courses/{course_id}/key-terms")
async def list_key_terms(course_id: str, user: dict = Depends(get_current_user)):
    items = await db.key_terms.find({"course_id": course_id, "user_id": user["id"]}, {"_id": 0}).sort("created_at", -1).to_list(100)
    return items

@api_router.get("/key-terms/{ktid}")
async def get_key_terms(ktid: str, user: dict = Depends(get_current_user)):
    k = await db.key_terms.find_one({"id": ktid, "user_id": user["id"]}, {"_id": 0})
    if not k:
        raise HTTPException(status_code=404, detail="Key terms not found")
    return k

@api_router.delete("/key-terms/{ktid}")
async def delete_key_terms(ktid: str, user: dict = Depends(get_current_user)):
    await db.key_terms.delete_one({"id": ktid, "user_id": user["id"]})
    return {"ok": True}

# -------------------- Paper Grader --------------------
@api_router.post("/paper-grader/grade")
async def grade_paper(request: Request, user: dict = Depends(get_current_user)):
    form = await request.form()
    paper_file = form.get("paper")
    rubric_file = form.get("rubric")
    title = form.get("title", "My Paper")
    previous_id = form.get("previous_id")
    if not paper_file:
        raise HTTPException(status_code=400, detail="Paper file is required")
    paper_data = await paper_file.read()
    paper_text = extract_text(paper_data, paper_file.filename.split(".")[-1].lower() if "." in paper_file.filename else "txt")
    rubric_text = ""
    if rubric_file:
        rubric_data = await rubric_file.read()
        rubric_text = extract_text(rubric_data, rubric_file.filename.split(".")[-1].lower() if "." in rubric_file.filename else "txt")
    rubric_prompt = f"\n\nRubric:\n{rubric_text}" if rubric_text else "\n\nUse default academic grading criteria: Comprehension, Structure, Scientific Accuracy, Depth of Analysis. Score each out of 10."
    system = "You are an expert paper grader. Respond with VALID JSON only."
    prompt = f"""Grade the following academic paper. Provide a score out of 10 for each criterion, explain what the student did well, and provide specific improvements to reach 10/10.
Paper title: {title}
Paper content:
{paper_text[:80000]}
{rubric_prompt}
Format as JSON:
{{"overall_percentage": 85, "summary": "brief summary", "criteria": [{{"name": "Criterion name", "score": 8, "did_well": "what they did well", "improve": "how to reach 10/10"}}]}}
If any criterion is already 10/10, set improve to null."""
    try:
        raw = await llm_generate(system, prompt)
        parsed = parse_json_block(raw)
        result = sanitize_obj(parsed)
    except Exception as e:
        logger.error(f"paper grader failed: {e}")
        raise HTTPException(status_code=500, detail="Paper grading failed. Please try again.")
    attempt_num = 1
    if previous_id:
        prev = await db.paper_grades.find_one({"id": previous_id, "user_id": user["id"]})
        if prev:
            attempt_num = prev.get("attempt", 0) + 1
    doc = {"id": str(uuid.uuid4()), "user_id": user["id"], "title": title, "paper_text": paper_text[:10000], "rubric_text": rubric_text[:5000] if rubric_text else None, "rubric_provided": bool(rubric_text), "attempt": attempt_num, "previous_id": previous_id, "overall_percentage": result.get("overall_percentage", 0), "perfect": result.get("overall_percentage", 0) == 100, "summary": result.get("summary", ""), "criteria": result.get("criteria", []), "created_at": now_iso()}
    await db.paper_grades.insert_one(doc)
    doc.pop("_id", None)
    return doc

@api_router.get("/paper-grader")
async def list_grades(user: dict = Depends(get_current_user)):
    grades = await db.paper_grades.find({"user_id": user["id"]}, {"_id": 0}).sort("created_at", -1).to_list(100)
    return grades

@api_router.get("/paper-grader/{gid}")
async def get_grade(gid: str, user: dict = Depends(get_current_user)):
    g = await db.paper_grades.find_one({"id": gid, "user_id": user["id"]}, {"_id": 0})
    if not g:
        raise HTTPException(status_code=404, detail="Grade not found")
    return g

@api_router.delete("/paper-grader/{gid}")
async def delete_grade(gid: str, user: dict = Depends(get_current_user)):
    await db.paper_grades.delete_one({"id": gid, "user_id": user["id"]})
    return {"ok": True}

# Include router and health route
app.include_router(api_router)

@app.get("/health")
async def health_check():
    return {"status": "ok"}

@app.get("/")
async def root():
    return {"message": "StudyAce API is running"}
